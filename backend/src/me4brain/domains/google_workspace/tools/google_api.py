"""Google Workspace API Tools.

Questo modulo contiene wrapper async per GoogleWorkspaceService.
I tool vengono chiamati dal GoogleWorkspaceHandler.

Tools:
- drive_search: Cerca file su Google Drive
- drive_list_files: Elenca file in cartella
- drive_get_file: Ottiene metadati file
- gmail_search: Cerca email
- gmail_get_message: Legge email
- calendar_upcoming: Eventi prossimi giorni
- calendar_list_events: Lista eventi con filtro
- calendar_analyze_meetings: Analizza meeting da tutte le piattaforme
"""

import asyncio
from concurrent.futures import ThreadPoolExecutor
from datetime import UTC
from functools import partial
from typing import Any

import structlog

# Import calendar meeting analyzer
from me4brain.domains.google_workspace.tools.calendar_meeting_analyzer import (
    calendar_analyze_meetings,
)

logger = structlog.get_logger(__name__)

# Thread pool per chiamate sync
_executor = ThreadPoolExecutor(max_workers=4)


async def _run_sync(func, *args, **kwargs) -> Any:
    """Esegue funzione sync in thread pool."""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(
        _executor,
        partial(func, *args, **kwargs),
    )


def _get_workspace_service():
    """Ottiene istanza GoogleWorkspaceService."""
    from me4brain.integrations.google_workspace import get_google_workspace_service

    return get_google_workspace_service()


async def _extract_text_from_binary(file_bytes: bytes, mime_type: str, file_name: str) -> str:
    """Extract text from binary Office files (docx, xlsx, pptx) or PDF.

    Args:
        file_bytes: Raw file bytes
        mime_type: MIME type of the file
        file_name: Original file name

    Returns:
        Extracted text content
    """
    import io

    try:
        # DOCX - Word documents
        if "wordprocessingml" in mime_type or file_name.endswith(".docx"):
            try:
                from docx import Document

                doc = Document(io.BytesIO(file_bytes))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return "\n\n".join(paragraphs)
            except ImportError:
                logger.warning("docx_library_not_installed")
                return "[DOCX extraction requires python-docx library]"
            except Exception as e:
                return f"[Error extracting DOCX: {e}]"

        # XLSX - Excel spreadsheets
        elif (
            "spreadsheetml" in mime_type
            or "ms-excel" in mime_type
            or file_name.endswith(".xlsx")
            or file_name.endswith(".xls")
        ):
            try:
                import pandas as pd

                # Read all sheets
                xlsx = pd.ExcelFile(io.BytesIO(file_bytes))
                all_text = []
                for sheet_name in xlsx.sheet_names[:5]:  # Limit to first 5 sheets
                    df = pd.read_excel(xlsx, sheet_name=sheet_name)
                    all_text.append(f"### Sheet: {sheet_name}\n{df.to_string()}")
                return "\n\n".join(all_text)
            except ImportError:
                logger.warning("pandas_openpyxl_not_installed")
                return "[XLSX extraction requires pandas and openpyxl libraries]"
            except Exception as e:
                return f"[Error extracting XLSX: {e}]"

        # PPTX - PowerPoint presentations
        elif (
            "presentationml" in mime_type
            or "ms-powerpoint" in mime_type
            or file_name.endswith(".pptx")
        ):
            try:
                from pptx import Presentation

                prs = Presentation(io.BytesIO(file_bytes))
                texts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_texts.append(shape.text)
                    if slide_texts:
                        texts.append(f"### Slide {slide_num}\n" + "\n".join(slide_texts))
                return "\n\n".join(texts)
            except ImportError:
                logger.warning("pptx_library_not_installed")
                return "[PPTX extraction requires python-pptx library]"
            except Exception as e:
                return f"[Error extracting PPTX: {e}]"

        # PDF
        elif mime_type == "application/pdf" or file_name.endswith(".pdf"):
            try:
                import pymupdf  # PyMuPDF

                doc = pymupdf.open(stream=file_bytes, filetype="pdf")
                texts = []
                for page_num, page in enumerate(doc, 1):
                    text = page.get_text()
                    if text.strip():
                        texts.append(f"### Page {page_num}\n{text}")
                doc.close()
                return "\n\n".join(texts)
            except ImportError:
                logger.warning("pymupdf_not_installed")
                return "[PDF extraction requires pymupdf library]"
            except Exception as e:
                return f"[Error extracting PDF: {e}]"

        return f"[Unsupported file type: {mime_type}]"

    except Exception as e:
        logger.error("text_extraction_error", error=str(e), mime_type=mime_type)
        return f"[Error extracting text: {e}]"


# =============================================================================
# Drive Tools
# =============================================================================


async def drive_search(
    query: str,
    max_results: int = 20,
    mime_type: str | None = None,
    folder_id: str | None = None,
) -> dict[str, Any]:
    """Cerca file su Google Drive.

    Args:
        query: Termine di ricerca (full-text search)
        max_results: Numero massimo risultati
        mime_type: Filtra per tipo MIME (opzionale)
        folder_id: ID cartella per restringere la ricerca (opzionale)

    Returns:
        dict con lista file trovati
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Drive",
            }

        # Usa metodo search già implementato con folder_id
        result = await service.search(query=query, max_results=max_results, folder_id=folder_id)

        return {
            "query": query,
            "folder_id": folder_id,
            "files": result.get("files", []),
            "count": result.get("count", 0),
            "source": "Google Drive",
        }

    except Exception as e:
        logger.error("drive_search_error", error=str(e))
        return {"error": str(e), "source": "Google Drive"}


async def drive_list_files(
    folder_id: str | None = None,
    max_results: int = 50,
) -> dict[str, Any]:
    """Elenca file in una cartella Drive.

    Args:
        folder_id: ID cartella (None = root)
        max_results: Numero massimo risultati

    Returns:
        dict con lista file
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Drive",
            }

        result = await service.drive_list_files(
            folder_id=folder_id,
            max_results=max_results,
        )

        # Result è una lista o un dict con chiave files
        files = result.get("files", []) if isinstance(result, dict) else result

        return {
            "folder_id": folder_id or "root",
            "files": files,
            "count": len(files) if files else 0,
            "source": "Google Drive",
        }

    except Exception as e:
        logger.error("drive_list_files_error", error=str(e))
        return {"error": str(e), "source": "Google Drive"}


async def drive_get_file(file_id: str) -> dict[str, Any]:
    """Ottiene metadati di un file.

    Args:
        file_id: ID del file Drive

    Returns:
        dict con metadati file
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Drive",
            }

        file_info = await service.drive_get_file(file_id=file_id)

        return {
            "file": file_info,
            "source": "Google Drive",
        }

    except Exception as e:
        logger.error("drive_get_file_error", error=str(e))
        return {"error": str(e), "source": "Google Drive"}


async def drive_get_content(file_id: str) -> dict[str, Any]:
    """Estrae il contenuto testuale di un file Google (Docs, Sheets, Slides).

    Args:
        file_id: ID del file Drive

    Returns:
        dict con contenuto testuale estratto
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Drive",
            }

        drive = await _run_sync(partial(service._get_google_service, "drive", "v3"))

        # Prima ottieni metadati per capire il tipo
        file_meta = await _run_sync(
            lambda: drive.files().get(fileId=file_id, fields="id,name,mimeType").execute()
        )

        mime_type = file_meta.get("mimeType", "")
        file_name = file_meta.get("name", "")

        # Export based on MIME type
        export_mime = None
        is_office_file = False
        is_pdf = False

        # Google native files - can export
        if "document" in mime_type and "google-apps" in mime_type:
            export_mime = "text/plain"
        elif "spreadsheet" in mime_type and "google-apps" in mime_type:
            export_mime = "text/csv"
        elif "presentation" in mime_type and "google-apps" in mime_type:
            export_mime = "text/plain"
        # Office files - need to download and parse
        elif (
            "wordprocessingml" in mime_type
            or ".document" in mime_type
            or "spreadsheetml" in mime_type
            or "ms-excel" in mime_type
            or "presentationml" in mime_type
            or "ms-powerpoint" in mime_type
        ):
            is_office_file = True
        elif mime_type == "application/pdf" or file_name.endswith(".pdf"):
            is_pdf = True

        content = ""

        if export_mime:
            # Google native files - export
            content = await _run_sync(
                lambda: drive.files().export(fileId=file_id, mimeType=export_mime).execute()
            )
            if isinstance(content, bytes):
                content = content.decode("utf-8", errors="replace")

        elif is_office_file or is_pdf:
            # Download file and extract text

            file_bytes = await _run_sync(lambda: drive.files().get_media(fileId=file_id).execute())
            if isinstance(file_bytes, bytes):
                content = await _extract_text_from_binary(file_bytes, mime_type, file_name)
        else:
            # Unknown binary file
            file_bytes = await _run_sync(lambda: drive.files().get_media(fileId=file_id).execute())
            if isinstance(file_bytes, bytes):
                content = f"[Binary file: {len(file_bytes)} bytes - content extraction not supported for this type]"

        return {
            "file_id": file_id,
            "name": file_name,
            "mime_type": mime_type,
            "content": content[:50000] if content else "",  # Limita a 50k chars
            "source": "Google Drive",
        }

    except Exception as e:
        logger.error("drive_get_content_error", error=str(e))
        return {"error": str(e), "source": "Google Drive"}


async def drive_export(
    file_id: str,
    export_format: str = "pdf",
) -> dict[str, Any]:
    """Esporta file Google in formato standard.

    Args:
        file_id: ID del file
        export_format: Formato output (pdf, docx, xlsx, pptx, txt, csv)

    Returns:
        dict con info esportazione (base64 content)
    """
    import base64

    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Drive",
            }

        drive = await _run_sync(partial(service._get_google_service, "drive", "v3"))

        # Mappa formati
        format_map = {
            "pdf": "application/pdf",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "txt": "text/plain",
            "csv": "text/csv",
        }

        export_mime = format_map.get(export_format.lower(), "application/pdf")

        content = await _run_sync(
            lambda: drive.files().export(fileId=file_id, mimeType=export_mime).execute()
        )

        # Encode in base64 per trasporto
        if isinstance(content, bytes):
            content_b64 = base64.b64encode(content).decode("utf-8")
            size_bytes = len(content)
        else:
            content_b64 = base64.b64encode(content.encode()).decode("utf-8")
            size_bytes = len(content)

        return {
            "file_id": file_id,
            "format": export_format,
            "mime_type": export_mime,
            "size_bytes": size_bytes,
            "content_base64": content_b64[:10000] + "..."
            if len(content_b64) > 10000
            else content_b64,
            "source": "Google Drive",
        }

    except Exception as e:
        logger.error("drive_export_error", error=str(e))
        return {"error": str(e), "source": "Google Drive"}


async def drive_create_folder(
    name: str,
    parent_id: str | None = None,
) -> dict[str, Any]:
    """Crea una nuova cartella su Drive.

    Args:
        name: Nome cartella
        parent_id: ID cartella parent (None = root)

    Returns:
        dict con info cartella creata
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Drive",
            }

        drive = await _run_sync(partial(service._get_google_service, "drive", "v3"))

        file_metadata = {
            "name": name,
            "mimeType": "application/vnd.google-apps.folder",
        }

        if parent_id:
            file_metadata["parents"] = [parent_id]

        folder = await _run_sync(
            lambda: drive.files().create(body=file_metadata, fields="id,name,webViewLink").execute()
        )

        return {
            "folder_id": folder.get("id"),
            "name": folder.get("name"),
            "link": folder.get("webViewLink"),
            "source": "Google Drive",
        }

    except Exception as e:
        logger.error("drive_create_folder_error", error=str(e))
        return {"error": str(e), "source": "Google Drive"}


async def drive_copy(
    file_id: str,
    new_name: str | None = None,
    parent_id: str | None = None,
) -> dict[str, Any]:
    """Copia un file esistente.

    Args:
        file_id: ID file da copiare
        new_name: Nome nuovo file (opzionale)
        parent_id: ID cartella destinazione (opzionale)

    Returns:
        dict con info file copiato
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Drive",
            }

        drive = await _run_sync(partial(service._get_google_service, "drive", "v3"))

        body = {}
        if new_name:
            body["name"] = new_name
        if parent_id:
            body["parents"] = [parent_id]

        copied = await _run_sync(
            lambda: (
                drive.files()
                .copy(
                    fileId=file_id,
                    body=body if body else None,
                    fields="id,name,webViewLink",
                )
                .execute()
            )
        )

        return {
            "original_id": file_id,
            "new_id": copied.get("id"),
            "name": copied.get("name"),
            "link": copied.get("webViewLink"),
            "source": "Google Drive",
        }

    except Exception as e:
        logger.error("drive_copy_error", error=str(e))
        return {"error": str(e), "source": "Google Drive"}


# =============================================================================
# Gmail Tools
# =============================================================================


async def gmail_search(
    query: str,
    max_results: int = 20,
) -> dict[str, Any]:
    """Cerca email in Gmail.

    Args:
        query: REQUIRED. Stringa di ricerca per Gmail. Può essere una parola chiave (es. "ANCI") o una query complessa (es. "from:user@example.com after:2024/01/01").
        max_results: Numero massimo di email da restituire (default 20).

    Returns:
        dict con lista email
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Gmail"}

        emails = await service.gmail_search(query=query, max_results=max_results)

        # Format emails per output
        email_summaries = []
        for email in emails:
            # Estrai header per from/subject/date
            headers = {h["name"]: h["value"] for h in email.get("payload", {}).get("headers", [])}
            email_summaries.append(
                {
                    "id": email.get("id"),
                    "subject": headers.get("Subject", "(no subject)"),
                    "from": headers.get("From"),
                    "date": headers.get("Date"),
                    "snippet": email.get("snippet", "")[:200],
                }
            )

        return {
            "query": query,
            "emails": email_summaries,
            "count": len(email_summaries),
            "source": "Gmail",
        }

    except Exception as e:
        logger.error("gmail_search_error", error=str(e))
        return {"error": str(e), "source": "Gmail"}


async def gmail_get_message(message_id: str) -> dict[str, Any]:
    """Legge una email specifica.

    Args:
        message_id: ID del messaggio Gmail

    Returns:
        dict con contenuto email
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Gmail"}

        message = await service.gmail_get_message(message_id=message_id)

        return {
            "message": message,
            "source": "Gmail",
        }

    except Exception as e:
        logger.error("gmail_get_message_error", error=str(e))
        return {"error": str(e), "source": "Gmail"}


async def gmail_send(
    to: str,
    subject: str,
    body: str,
    cc: str | None = None,
    bcc: str | None = None,
) -> dict[str, Any]:
    """Invia una nuova email.

    Args:
        to: Destinatario
        subject: Oggetto
        body: Corpo email (text/plain)
        cc: CC (opzionale)
        bcc: BCC (opzionale)

    Returns:
        dict con info email inviata
    """
    import base64
    from email.mime.text import MIMEText

    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Gmail"}

        gmail = await _run_sync(partial(service._get_google_service, "gmail", "v1"))

        message = MIMEText(body)
        message["to"] = to
        message["subject"] = subject

        if cc:
            message["cc"] = cc
        if bcc:
            message["bcc"] = bcc

        raw = base64.urlsafe_b64encode(message.as_bytes()).decode("utf-8")

        result = await _run_sync(
            lambda: gmail.users().messages().send(userId="me", body={"raw": raw}).execute()
        )

        return {
            "message_id": result.get("id"),
            "thread_id": result.get("threadId"),
            "to": to,
            "subject": subject,
            "source": "Gmail",
        }

    except Exception as e:
        logger.error("gmail_send_error", error=str(e))
        return {"error": str(e), "source": "Gmail"}


async def gmail_reply(
    message_id: str,
    body: str,
) -> dict[str, Any]:
    """Risponde a una email esistente.

    Args:
        message_id: ID messaggio originale
        body: Corpo risposta

    Returns:
        dict con info risposta inviata
    """
    import base64
    from email.mime.text import MIMEText

    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Gmail"}

        gmail = await _run_sync(partial(service._get_google_service, "gmail", "v1"))

        # Ottieni messaggio originale per thread_id e headers
        original = await _run_sync(
            lambda: (
                gmail.users()
                .messages()
                .get(
                    userId="me",
                    id=message_id,
                    format="metadata",
                    metadataHeaders=["Subject", "From"],
                )
                .execute()
            )
        )

        thread_id = original.get("threadId")
        headers = {h["name"]: h["value"] for h in original.get("payload", {}).get("headers", [])}

        original_subject = headers.get("Subject", "")
        reply_to = headers.get("From", "")

        # Prepara reply
        reply_subject = (
            original_subject if original_subject.startswith("Re:") else f"Re: {original_subject}"
        )

        message = MIMEText(body)
        message["to"] = reply_to
        message["subject"] = reply_subject
        message["In-Reply-To"] = message_id
        message["References"] = message_id

        raw = base64.urlsafe_b64encode(message.as_bytes()).decode("utf-8")

        result = await _run_sync(
            lambda: (
                gmail.users()
                .messages()
                .send(userId="me", body={"raw": raw, "threadId": thread_id})
                .execute()
            )
        )

        return {
            "message_id": result.get("id"),
            "thread_id": result.get("threadId"),
            "reply_to": reply_to,
            "subject": reply_subject,
            "source": "Gmail",
        }

    except Exception as e:
        logger.error("gmail_reply_error", error=str(e))
        return {"error": str(e), "source": "Gmail"}


async def gmail_forward(
    message_id: str,
    to: str,
    additional_text: str | None = None,
) -> dict[str, Any]:
    """Inoltra una email.

    Args:
        message_id: ID messaggio da inoltrare
        to: Destinatario forward
        additional_text: Testo aggiuntivo (opzionale)

    Returns:
        dict con info forward inviato
    """
    import base64
    from email.mime.text import MIMEText

    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Gmail"}

        gmail = await _run_sync(partial(service._get_google_service, "gmail", "v1"))

        # Ottieni messaggio originale
        original = await _run_sync(
            lambda: (
                gmail.users().messages().get(userId="me", id=message_id, format="full").execute()
            )
        )

        headers = {h["name"]: h["value"] for h in original.get("payload", {}).get("headers", [])}
        original_subject = headers.get("Subject", "")
        original_from = headers.get("From", "")
        original_date = headers.get("Date", "")

        # Estrai body originale
        snippet = original.get("snippet", "")

        # Componi forward
        forward_subject = (
            f"Fwd: {original_subject}"
            if not original_subject.startswith("Fwd:")
            else original_subject
        )

        forward_body = ""
        if additional_text:
            forward_body = f"{additional_text}\n\n"
        forward_body += "---------- Forwarded message ---------\n"
        forward_body += f"From: {original_from}\n"
        forward_body += f"Date: {original_date}\n"
        forward_body += f"Subject: {original_subject}\n\n"
        forward_body += snippet

        message = MIMEText(forward_body)
        message["to"] = to
        message["subject"] = forward_subject

        raw = base64.urlsafe_b64encode(message.as_bytes()).decode("utf-8")

        result = await _run_sync(
            lambda: gmail.users().messages().send(userId="me", body={"raw": raw}).execute()
        )

        return {
            "message_id": result.get("id"),
            "forwarded_to": to,
            "subject": forward_subject,
            "source": "Gmail",
        }

    except Exception as e:
        logger.error("gmail_forward_error", error=str(e))
        return {"error": str(e), "source": "Gmail"}


async def gmail_get_attachments(message_id: str) -> dict[str, Any]:
    """Recupera gli allegati di una email.

    Args:
        message_id: ID del messaggio Gmail

    Returns:
        dict con lista allegati (filename, mimeType, size)
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Gmail"}

        attachments = await service.gmail_get_attachments(message_id=message_id)

        # Rimuovi data base64 per output leggero, mantieni solo metadata + attachmentId
        attachment_summaries = [
            {
                "attachmentId": att.get("attachmentId"),
                "filename": att.get("filename"),
                "mimeType": att.get("mimeType"),
                "size": att.get("size"),
            }
            for att in attachments
        ]

        return {
            "message_id": message_id,
            "attachments": attachment_summaries,
            "count": len(attachment_summaries),
            "source": "Gmail",
        }

    except Exception as e:
        logger.error("gmail_get_attachments_error", error=str(e))
        return {"error": str(e), "source": "Gmail"}


async def gmail_get_attachment_content(
    message_id: str,
    attachment_id: str,
    filename: str = "",
    mime_type: str = "",
) -> dict[str, Any]:
    """Scarica un allegato e ne estrae il contenuto testuale.

    Args:
        message_id: ID del messaggio Gmail
        attachment_id: ID dell'allegato (da gmail_get_attachments)
        filename: Nome file originale (opzionale, per inferenza mimeType)
        mime_type: Tipo MIME (opzionale)

    Returns:
        dict con filename, mimeType, e contenuto testuale estratto
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Gmail"}

        # Download attachment as bytes, passing filename for mimeType inference
        att_info = await service.gmail_get_attachment_content(
            message_id=message_id,
            attachment_id=attachment_id,
            filename=filename,
            mime_type=mime_type,
        )

        filename = att_info.get("filename", "")
        mime_type = att_info.get("mimeType", "")
        raw_bytes = att_info.get("data", b"")

        # Extract text content
        text_content = await _extract_text_from_binary(raw_bytes, mime_type, filename)

        return {
            "message_id": message_id,
            "attachment_id": attachment_id,
            "filename": filename,
            "mimeType": mime_type,
            "size": len(raw_bytes),
            "content": text_content[:50000] if text_content else "",  # Limit to 50k chars
            "source": "Gmail",
        }

    except Exception as e:
        logger.error("gmail_get_attachment_content_error", error=str(e))
        return {"error": str(e), "source": "Gmail"}


# =============================================================================
# Calendar Tools
# =============================================================================


async def calendar_upcoming(
    days: int = 7,
    calendar_id: str = "primary",
) -> dict[str, Any]:
    """Ottiene eventi calendario prossimi N giorni.

    Args:
        days: Numero giorni da cercare
        calendar_id: ID calendario (default: primary)

    Returns:
        dict con lista eventi
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Calendar",
            }

        events = await service.calendar_get_upcoming(days=days, calendar_id=calendar_id)

        # Format eventi
        event_summaries = []
        for event in events:
            event_summaries.append(
                {
                    "id": event.get("id"),
                    "summary": event.get("summary", "(no title)"),
                    "start": event.get("start", {}).get("dateTime")
                    or event.get("start", {}).get("date"),
                    "end": event.get("end", {}).get("dateTime") or event.get("end", {}).get("date"),
                    "location": event.get("location"),
                    "status": event.get("status"),
                }
            )

        return {
            "days": days,
            "events": event_summaries,
            "count": len(event_summaries),
            "source": "Google Calendar",
        }

    except Exception as e:
        logger.error("calendar_upcoming_error", error=str(e))
        return {"error": str(e), "source": "Google Calendar"}


async def calendar_list_events(
    query: str | None = None,
    time_min: str | None = None,
    time_max: str | None = None,
    max_results: int = 50,
) -> dict[str, Any]:
    """Lista eventi calendario con filtri.

    Args:
        query: Termine di ricerca
        time_min: Data/ora minima (ISO format)
        time_max: Data/ora massima (ISO format)
        max_results: Numero massimo risultati

    Returns:
        dict con lista eventi
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Calendar",
            }

        # Only default time_min for forward-looking queries (no query specified)
        # For search queries with explicit keywords, don't filter by time unless requested
        if not time_min and not query:
            from datetime import UTC, datetime

            time_min = datetime.now(UTC).isoformat()

        events = await service.calendar_list_events(
            time_min=time_min,
            time_max=time_max,
            max_results=max_results,
            query=query,
        )

        return {
            "query": query,
            "events": events,
            "count": len(events) if events else 0,
            "source": "Google Calendar",
        }

    except Exception as e:
        logger.error("calendar_list_events_error", error=str(e))
        return {"error": str(e), "source": "Google Calendar"}


async def calendar_create_event(
    summary: str,
    start_time: str,
    end_time: str,
    description: str | None = None,
    location: str | None = None,
    attendees: list[str] | None = None,
    calendar_id: str = "primary",
) -> dict[str, Any]:
    """Crea un nuovo evento nel calendario.

    Args:
        summary: Titolo evento
        start_time: Inizio (ISO format, es. "2024-01-29T10:00:00+01:00")
        end_time: Fine (ISO format)
        description: Descrizione (opzionale)
        location: Luogo (opzionale)
        attendees: Lista email partecipanti (opzionale)
        calendar_id: ID calendario (default: primary)

    Returns:
        dict con info evento creato
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Calendar",
            }

        calendar = await _run_sync(partial(service._get_google_service, "calendar", "v3"))

        event = {
            "summary": summary,
            "start": {"dateTime": start_time, "timeZone": "Europe/Rome"},
            "end": {"dateTime": end_time, "timeZone": "Europe/Rome"},
        }

        if description:
            event["description"] = description
        if location:
            event["location"] = location
        if attendees:
            event["attendees"] = [{"email": email} for email in attendees]

        result = await _run_sync(
            lambda: calendar.events().insert(calendarId=calendar_id, body=event).execute()
        )

        return {
            "event_id": result.get("id"),
            "summary": result.get("summary"),
            "start": result.get("start", {}).get("dateTime"),
            "end": result.get("end", {}).get("dateTime"),
            "link": result.get("htmlLink"),
            "source": "Google Calendar",
        }

    except Exception as e:
        logger.error("calendar_create_event_error", error=str(e))
        return {"error": str(e), "source": "Google Calendar"}


async def calendar_get_event(
    event_id: str,
    calendar_id: str = "primary",
) -> dict[str, Any]:
    """Ottiene dettagli di un evento specifico.

    Args:
        event_id: ID evento
        calendar_id: ID calendario

    Returns:
        dict con dettagli evento
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Calendar",
            }

        calendar = await _run_sync(partial(service._get_google_service, "calendar", "v3"))

        event = await _run_sync(
            lambda: calendar.events().get(calendarId=calendar_id, eventId=event_id).execute()
        )

        return {
            "event_id": event.get("id"),
            "summary": event.get("summary"),
            "description": event.get("description"),
            "start": event.get("start", {}).get("dateTime") or event.get("start", {}).get("date"),
            "end": event.get("end", {}).get("dateTime") or event.get("end", {}).get("date"),
            "location": event.get("location"),
            "attendees": [a.get("email") for a in event.get("attendees", [])],
            "status": event.get("status"),
            "link": event.get("htmlLink"),
            "source": "Google Calendar",
        }

    except Exception as e:
        logger.error("calendar_get_event_error", error=str(e))
        return {"error": str(e), "source": "Google Calendar"}


async def calendar_update_event(
    event_id: str,
    summary: str | None = None,
    start_time: str | None = None,
    end_time: str | None = None,
    description: str | None = None,
    location: str | None = None,
    calendar_id: str = "primary",
) -> dict[str, Any]:
    """Aggiorna un evento esistente.

    Args:
        event_id: ID evento
        summary: Nuovo titolo (opzionale)
        start_time: Nuovo inizio (opzionale)
        end_time: Nuova fine (opzionale)
        description: Nuova descrizione (opzionale)
        location: Nuovo luogo (opzionale)
        calendar_id: ID calendario

    Returns:
        dict con info evento aggiornato
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Calendar",
            }

        calendar = await _run_sync(partial(service._get_google_service, "calendar", "v3"))

        # Prima ottieni evento esistente
        event = await _run_sync(
            lambda: calendar.events().get(calendarId=calendar_id, eventId=event_id).execute()
        )

        # Aggiorna campi specificati
        if summary:
            event["summary"] = summary
        if description:
            event["description"] = description
        if location:
            event["location"] = location
        if start_time:
            event["start"] = {"dateTime": start_time, "timeZone": "Europe/Rome"}
        if end_time:
            event["end"] = {"dateTime": end_time, "timeZone": "Europe/Rome"}

        result = await _run_sync(
            lambda: (
                calendar.events()
                .update(calendarId=calendar_id, eventId=event_id, body=event)
                .execute()
            )
        )

        return {
            "event_id": result.get("id"),
            "summary": result.get("summary"),
            "updated": result.get("updated"),
            "link": result.get("htmlLink"),
            "source": "Google Calendar",
        }

    except Exception as e:
        logger.error("calendar_update_event_error", error=str(e))
        return {"error": str(e), "source": "Google Calendar"}


async def calendar_delete_event(
    event_id: str,
    calendar_id: str = "primary",
) -> dict[str, Any]:
    """Cancella un evento dal calendario.

    Args:
        event_id: ID evento da cancellare
        calendar_id: ID calendario

    Returns:
        dict con conferma cancellazione
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Calendar",
            }

        calendar = await _run_sync(partial(service._get_google_service, "calendar", "v3"))

        await _run_sync(
            lambda: calendar.events().delete(calendarId=calendar_id, eventId=event_id).execute()
        )

        return {
            "event_id": event_id,
            "deleted": True,
            "source": "Google Calendar",
        }

    except Exception as e:
        logger.error("calendar_delete_event_error", error=str(e))
        return {"error": str(e), "source": "Google Calendar"}


# =============================================================================
# Docs Tools
# =============================================================================


async def docs_get(document_id: str) -> dict[str, Any]:
    """Ottiene contenuto di un documento Google Docs.

    Args:
        document_id: ID del documento

    Returns:
        dict con contenuto documento
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Google Docs"}

        doc = await _run_sync(partial(service._get_google_service, "docs", "v1"))

        result = await _run_sync(lambda: doc.documents().get(documentId=document_id).execute())

        # Estrai testo dal documento
        body = result.get("body", {})
        content_elements = body.get("content", [])
        text_content = []

        for element in content_elements:
            if "paragraph" in element:
                for elem in element["paragraph"].get("elements", []):
                    if "textRun" in elem:
                        text_content.append(elem["textRun"].get("content", ""))

        return {
            "document_id": document_id,
            "title": result.get("title", ""),
            "text": "".join(text_content),
            "source": "Google Docs",
        }

    except Exception as e:
        logger.error("docs_get_error", error=str(e))
        return {"error": str(e), "source": "Google Docs"}


async def docs_create(
    title: str,
    content: str | None = None,
    folder_id: str | None = None,
) -> dict[str, Any]:
    """Crea un nuovo documento Google Docs.

    Args:
        title: Titolo documento
        content: Contenuto iniziale (opzionale)
        folder_id: ID cartella Google Drive dove salvare il documento (opzionale)

    Returns:
        dict con info documento creato
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Google Docs"}

        docs = await _run_sync(partial(service._get_google_service, "docs", "v1"))

        # Crea documento vuoto
        doc = await _run_sync(lambda: docs.documents().create(body={"title": title}).execute())

        document_id = doc.get("documentId")

        # Se c'è contenuto, inseriscilo
        if content:
            await _run_sync(
                lambda: (
                    docs.documents()
                    .batchUpdate(
                        documentId=document_id,
                        body={
                            "requests": [
                                {
                                    "insertText": {
                                        "location": {"index": 1},
                                        "text": content,
                                    }
                                }
                            ]
                        },
                    )
                    .execute()
                )
            )

        # Se specificata una cartella, sposta il documento
        folder_link = None
        if folder_id:
            drive = await _run_sync(partial(service._get_google_service, "drive", "v3"))

            # Sposta il file nella cartella specificata
            # Prima ottieni il parent corrente (root)
            file_info = await _run_sync(
                lambda: drive.files().get(fileId=document_id, fields="parents").execute()
            )
            previous_parents = ",".join(file_info.get("parents", []))

            # Sposta nella nuova cartella
            await _run_sync(
                lambda: (
                    drive.files()
                    .update(
                        fileId=document_id,
                        addParents=folder_id,
                        removeParents=previous_parents,
                        fields="id, parents",
                    )
                    .execute()
                )
            )

            folder_link = f"https://drive.google.com/drive/folders/{folder_id}"
            logger.info("docs_create_moved_to_folder", document_id=document_id, folder_id=folder_id)

        result = {
            "document_id": document_id,
            "title": doc.get("title"),
            "link": f"https://docs.google.com/document/d/{document_id}/edit",
            "source": "Google Docs",
        }

        if folder_id:
            result["folder_id"] = folder_id
            result["folder_link"] = folder_link
            result["saved_to_folder"] = True

        return result

    except Exception as e:
        logger.error("docs_create_error", error=str(e))
        return {"error": str(e), "source": "Google Docs"}


async def docs_insert_text(
    document_id: str,
    text: str,
    index: int = 1,
) -> dict[str, Any]:
    """Inserisce testo in una posizione specifica del documento.

    Args:
        document_id: ID documento
        text: Testo da inserire
        index: Posizione (1 = inizio documento)

    Returns:
        dict con conferma inserimento
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Google Docs"}

        docs = await _run_sync(partial(service._get_google_service, "docs", "v1"))

        await _run_sync(
            lambda: (
                docs.documents()
                .batchUpdate(
                    documentId=document_id,
                    body={
                        "requests": [
                            {
                                "insertText": {
                                    "location": {"index": index},
                                    "text": text,
                                }
                            }
                        ]
                    },
                )
                .execute()
            )
        )

        return {
            "document_id": document_id,
            "inserted_length": len(text),
            "at_index": index,
            "source": "Google Docs",
        }

    except Exception as e:
        logger.error("docs_insert_text_error", error=str(e))
        return {"error": str(e), "source": "Google Docs"}


async def docs_append_text(
    document_id: str,
    text: str,
) -> dict[str, Any]:
    """Aggiunge testo alla fine del documento.

    Args:
        document_id: ID documento
        text: Testo da aggiungere

    Returns:
        dict con conferma aggiunta
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Google Docs"}

        docs = await _run_sync(partial(service._get_google_service, "docs", "v1"))

        # Prima ottieni lunghezza documento
        doc = await _run_sync(lambda: docs.documents().get(documentId=document_id).execute())

        body = doc.get("body", {})
        content = body.get("content", [])
        # L'ultimo elemento contiene l'endIndex
        end_index = content[-1].get("endIndex", 1) if content else 1

        await _run_sync(
            lambda: (
                docs.documents()
                .batchUpdate(
                    documentId=document_id,
                    body={
                        "requests": [
                            {
                                "insertText": {
                                    "location": {"index": end_index - 1},
                                    "text": text,
                                }
                            }
                        ]
                    },
                )
                .execute()
            )
        )

        return {
            "document_id": document_id,
            "appended_length": len(text),
            "source": "Google Docs",
        }

    except Exception as e:
        logger.error("docs_append_text_error", error=str(e))
        return {"error": str(e), "source": "Google Docs"}


async def docs_replace_text(
    document_id: str,
    find_text: str,
    replace_text: str,
    match_case: bool = False,
) -> dict[str, Any]:
    """Sostituisce tutte le occorrenze di testo nel documento.

    Args:
        document_id: ID documento
        find_text: Testo da cercare
        replace_text: Testo sostitutivo
        match_case: Case sensitive (default: False)

    Returns:
        dict con numero sostituzioni effettuate
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Google Docs"}

        docs = await _run_sync(partial(service._get_google_service, "docs", "v1"))

        result = await _run_sync(
            lambda: (
                docs.documents()
                .batchUpdate(
                    documentId=document_id,
                    body={
                        "requests": [
                            {
                                "replaceAllText": {
                                    "containsText": {
                                        "text": find_text,
                                        "matchCase": match_case,
                                    },
                                    "replaceText": replace_text,
                                }
                            }
                        ]
                    },
                )
                .execute()
            )
        )

        # Conta sostituzioni
        replies = result.get("replies", [])
        occurrences = (
            replies[0].get("replaceAllText", {}).get("occurrencesChanged", 0) if replies else 0
        )

        return {
            "document_id": document_id,
            "find_text": find_text,
            "replace_text": replace_text,
            "occurrences_replaced": occurrences,
            "source": "Google Docs",
        }

    except Exception as e:
        logger.error("docs_replace_text_error", error=str(e))
        return {"error": str(e), "source": "Google Docs"}


# =============================================================================
# Sheets Tools
# =============================================================================


async def sheets_get_values(
    spreadsheet_id: str,
    range_notation: str = "A1:Z100",
) -> dict[str, Any]:
    """Legge valori da un Google Sheet.

    Args:
        spreadsheet_id: ID dello spreadsheet
        range_notation: Range di celle (es. "Sheet1!A1:D10")

    Returns:
        dict con valori letti
    """
    try:
        service = _get_workspace_service()
        sheets = await _run_sync(partial(service._get_google_service, "sheets", "v4"))

        result = await _run_sync(
            lambda: (
                sheets.spreadsheets()
                .values()
                .get(
                    spreadsheetId=spreadsheet_id,
                    range=range_notation,
                )
                .execute()
            )
        )

        return {
            "spreadsheet_id": spreadsheet_id,
            "range": result.get("range", range_notation),
            "values": result.get("values", []),
            "source": "Google Sheets",
        }

    except Exception as e:
        logger.error("sheets_get_values_error", error=str(e))
        return {"error": str(e), "source": "Google Sheets"}


async def sheets_get_metadata(spreadsheet_id: str) -> dict[str, Any]:
    """Ottiene metadati di uno spreadsheet.

    Args:
        spreadsheet_id: ID dello spreadsheet

    Returns:
        dict con metadati spreadsheet
    """
    try:
        service = _get_workspace_service()
        sheets = await _run_sync(partial(service._get_google_service, "sheets", "v4"))

        result = await _run_sync(
            lambda: (
                sheets.spreadsheets()
                .get(
                    spreadsheetId=spreadsheet_id,
                )
                .execute()
            )
        )

        sheet_names = [s.get("properties", {}).get("title") for s in result.get("sheets", [])]

        return {
            "spreadsheet_id": spreadsheet_id,
            "title": result.get("properties", {}).get("title", ""),
            "sheets": sheet_names,
            "source": "Google Sheets",
        }

    except Exception as e:
        logger.error("sheets_get_metadata_error", error=str(e))
        return {"error": str(e), "source": "Google Sheets"}


async def sheets_create(
    title: str,
    sheet_names: list[str] | None = None,
) -> dict[str, Any]:
    """Crea un nuovo Google Spreadsheet.

    Args:
        title: Titolo spreadsheet
        sheet_names: Nomi fogli da creare (opzionale)

    Returns:
        dict con info spreadsheet creato
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Sheets",
            }

        sheets = await _run_sync(partial(service._get_google_service, "sheets", "v4"))

        body = {"properties": {"title": title}}

        if sheet_names:
            body["sheets"] = [{"properties": {"title": name}} for name in sheet_names]

        result = await _run_sync(lambda: sheets.spreadsheets().create(body=body).execute())

        return {
            "spreadsheet_id": result.get("spreadsheetId"),
            "title": result.get("properties", {}).get("title"),
            "sheets": [s.get("properties", {}).get("title") for s in result.get("sheets", [])],
            "link": result.get("spreadsheetUrl"),
            "source": "Google Sheets",
        }

    except Exception as e:
        logger.error("sheets_create_error", error=str(e))
        return {"error": str(e), "source": "Google Sheets"}


async def sheets_update_values(
    spreadsheet_id: str,
    range_notation: str,
    values: list[list[str]],
) -> dict[str, Any]:
    """Scrive valori in un range di celle.

    Args:
        spreadsheet_id: ID spreadsheet
        range_notation: Range (es. "Sheet1!A1:B2")
        values: Matrice valori [[r1c1, r1c2], [r2c1, r2c2]]

    Returns:
        dict con conferma scrittura
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Sheets",
            }

        sheets = await _run_sync(partial(service._get_google_service, "sheets", "v4"))

        result = await _run_sync(
            lambda: (
                sheets.spreadsheets()
                .values()
                .update(
                    spreadsheetId=spreadsheet_id,
                    range=range_notation,
                    valueInputOption="USER_ENTERED",
                    body={"values": values},
                )
                .execute()
            )
        )

        return {
            "spreadsheet_id": spreadsheet_id,
            "range": result.get("updatedRange"),
            "rows_updated": result.get("updatedRows", 0),
            "cells_updated": result.get("updatedCells", 0),
            "source": "Google Sheets",
        }

    except Exception as e:
        logger.error("sheets_update_values_error", error=str(e))
        return {"error": str(e), "source": "Google Sheets"}


async def sheets_append_row(
    spreadsheet_id: str,
    values: list[str],
    sheet_name: str = "Sheet1",
) -> dict[str, Any]:
    """Aggiunge una riga alla fine del foglio.

    Args:
        spreadsheet_id: ID spreadsheet
        values: Valori riga [col1, col2, col3, ...]
        sheet_name: Nome foglio

    Returns:
        dict con conferma aggiunta
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Sheets",
            }

        sheets = await _run_sync(partial(service._get_google_service, "sheets", "v4"))

        result = await _run_sync(
            lambda: (
                sheets.spreadsheets()
                .values()
                .append(
                    spreadsheetId=spreadsheet_id,
                    range="A1",  # Simple A1 works with default sheet
                    valueInputOption="USER_ENTERED",
                    insertDataOption="INSERT_ROWS",
                    body={"values": [values]},
                )
                .execute()
            )
        )

        return {
            "spreadsheet_id": spreadsheet_id,
            "range": result.get("updates", {}).get("updatedRange"),
            "rows_appended": result.get("updates", {}).get("updatedRows", 0),
            "source": "Google Sheets",
        }

    except Exception as e:
        logger.error("sheets_append_row_error", error=str(e))
        return {"error": str(e), "source": "Google Sheets"}


async def sheets_add_sheet(
    spreadsheet_id: str,
    sheet_title: str,
) -> dict[str, Any]:
    """Aggiunge un nuovo foglio allo spreadsheet.

    Args:
        spreadsheet_id: ID spreadsheet
        sheet_title: Nome nuovo foglio

    Returns:
        dict con info foglio creato
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Sheets",
            }

        sheets = await _run_sync(partial(service._get_google_service, "sheets", "v4"))

        result = await _run_sync(
            lambda: (
                sheets.spreadsheets()
                .batchUpdate(
                    spreadsheetId=spreadsheet_id,
                    body={"requests": [{"addSheet": {"properties": {"title": sheet_title}}}]},
                )
                .execute()
            )
        )

        replies = result.get("replies", [])
        sheet_props = replies[0].get("addSheet", {}).get("properties", {}) if replies else {}

        return {
            "spreadsheet_id": spreadsheet_id,
            "sheet_id": sheet_props.get("sheetId"),
            "sheet_title": sheet_props.get("title"),
            "source": "Google Sheets",
        }

    except Exception as e:
        logger.error("sheets_add_sheet_error", error=str(e))
        return {"error": str(e), "source": "Google Sheets"}


# =============================================================================
# Slides Tools
# =============================================================================


async def slides_get(presentation_id: str) -> dict[str, Any]:
    """Ottiene info su una presentazione Google Slides.

    Args:
        presentation_id: ID della presentazione

    Returns:
        dict con info presentazione
    """
    try:
        service = _get_workspace_service()
        slides = await _run_sync(partial(service._get_google_service, "slides", "v1"))

        result = await _run_sync(
            lambda: (
                slides.presentations()
                .get(
                    presentationId=presentation_id,
                )
                .execute()
            )
        )

        slide_count = len(result.get("slides", []))

        return {
            "presentation_id": presentation_id,
            "title": result.get("title", ""),
            "slide_count": slide_count,
            "source": "Google Slides",
        }

    except Exception as e:
        logger.error("slides_get_error", error=str(e))
        return {"error": str(e), "source": "Google Slides"}


async def slides_create(
    title: str,
) -> dict[str, Any]:
    """Crea una nuova presentazione Google Slides.

    Args:
        title: Titolo presentazione

    Returns:
        dict con info presentazione creata
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Slides",
            }

        slides = await _run_sync(partial(service._get_google_service, "slides", "v1"))

        result = await _run_sync(
            lambda: slides.presentations().create(body={"title": title}).execute()
        )

        return {
            "presentation_id": result.get("presentationId"),
            "title": result.get("title"),
            "slide_count": len(result.get("slides", [])),
            "link": f"https://docs.google.com/presentation/d/{result.get('presentationId')}/edit",
            "source": "Google Slides",
        }

    except Exception as e:
        logger.error("slides_create_error", error=str(e))
        return {"error": str(e), "source": "Google Slides"}


async def slides_get_text(
    presentation_id: str,
) -> dict[str, Any]:
    """Estrae tutto il testo da una presentazione.

    Args:
        presentation_id: ID presentazione

    Returns:
        dict con testo per ogni slide
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Slides",
            }

        slides = await _run_sync(partial(service._get_google_service, "slides", "v1"))

        result = await _run_sync(
            lambda: slides.presentations().get(presentationId=presentation_id).execute()
        )

        slides_text = []
        for idx, slide in enumerate(result.get("slides", []), 1):
            slide_texts = []
            for element in slide.get("pageElements", []):
                shape = element.get("shape", {})
                text_elements = shape.get("text", {}).get("textElements", [])
                for te in text_elements:
                    text_run = te.get("textRun", {})
                    content = text_run.get("content", "").strip()
                    if content:
                        slide_texts.append(content)

            slides_text.append(
                {
                    "slide_number": idx,
                    "text": " ".join(slide_texts),
                }
            )

        return {
            "presentation_id": presentation_id,
            "title": result.get("title", ""),
            "slide_count": len(result.get("slides", [])),
            "slides": slides_text,
            "source": "Google Slides",
        }

    except Exception as e:
        logger.error("slides_get_text_error", error=str(e))
        return {"error": str(e), "source": "Google Slides"}


async def slides_add_slide(
    presentation_id: str,
    layout: str = "BLANK",
) -> dict[str, Any]:
    """Aggiunge una nuova slide alla presentazione.

    Args:
        presentation_id: ID presentazione
        layout: Layout slide (BLANK, TITLE, TITLE_AND_BODY, TITLE_AND_TWO_COLUMNS, etc.)

    Returns:
        dict con info slide creata
    """
    import uuid

    try:
        service = _get_workspace_service()
        if not service:
            return {
                "error": "Google Workspace not configured",
                "source": "Google Slides",
            }

        slides = await _run_sync(partial(service._get_google_service, "slides", "v1"))

        # Genera ID unico per la nuova slide
        slide_id = f"slide_{uuid.uuid4().hex[:8]}"

        # Layout predefiniti
        layout_map = {
            "BLANK": "BLANK",
            "TITLE": "TITLE",
            "TITLE_AND_BODY": "TITLE_AND_BODY",
            "TITLE_AND_TWO_COLUMNS": "TITLE_AND_TWO_COLUMNS",
            "TITLE_ONLY": "TITLE_ONLY",
            "ONE_COLUMN_TEXT": "ONE_COLUMN_TEXT",
            "MAIN_POINT": "MAIN_POINT",
            "BIG_NUMBER": "BIG_NUMBER",
        }

        predefined_layout = layout_map.get(layout.upper(), "BLANK")

        await _run_sync(
            lambda: (
                slides.presentations()
                .batchUpdate(
                    presentationId=presentation_id,
                    body={
                        "requests": [
                            {
                                "createSlide": {
                                    "objectId": slide_id,
                                    "slideLayoutReference": {"predefinedLayout": predefined_layout},
                                }
                            }
                        ]
                    },
                )
                .execute()
            )
        )

        return {
            "presentation_id": presentation_id,
            "slide_id": slide_id,
            "layout": predefined_layout,
            "source": "Google Slides",
        }

    except Exception as e:
        logger.error("slides_add_slide_error", error=str(e))
        return {"error": str(e), "source": "Google Slides"}


# =============================================================================
# Meet Tools
# =============================================================================


async def meet_create(
    summary: str,
    duration_minutes: int = 60,
) -> dict[str, Any]:
    """Crea un meeting Google Meet via Calendar.

    Args:
        summary: Titolo del meeting
        duration_minutes: Durata in minuti

    Returns:
        dict con link meeting
    """
    try:
        from datetime import datetime, timedelta

        UTC = UTC

        service = _get_workspace_service()
        calendar = await _run_sync(partial(service._get_google_service, "calendar", "v3"))

        # Crea evento con conferenceData
        now = datetime.now(UTC)
        end_time = now + timedelta(minutes=duration_minutes)

        event = {
            "summary": summary,
            "start": {"dateTime": now.isoformat(), "timeZone": "UTC"},
            "end": {"dateTime": end_time.isoformat(), "timeZone": "UTC"},
            "conferenceData": {
                "createRequest": {
                    "requestId": f"meet-{now.timestamp()}",
                    "conferenceSolutionKey": {"type": "hangoutsMeet"},
                }
            },
        }

        result = await _run_sync(
            lambda: (
                calendar.events()
                .insert(
                    calendarId="primary",
                    body=event,
                    conferenceDataVersion=1,
                )
                .execute()
            )
        )

        meet_link = result.get("hangoutLink", "")

        return {
            "event_id": result.get("id", ""),
            "summary": summary,
            "meet_link": meet_link,
            "source": "Google Meet",
        }

    except Exception as e:
        logger.error("meet_create_error", error=str(e))
        return {"error": str(e), "source": "Google Meet"}


async def meet_list_conferences(
    max_results: int = 50,
) -> dict[str, Any]:
    """Lista le conferenze Google Meet passate.

    Args:
        max_results: Numero massimo di conferenze da restituire

    Returns:
        dict con lista conferenze e partecipanti
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Google Meet"}

        conferences = await service.meet_list_conferences(max_results=max_results)

        return {
            "conferences": conferences,
            "count": len(conferences),
            "source": "Google Meet",
        }

    except Exception as e:
        logger.error("meet_list_conferences_error", error=str(e))
        return {"error": str(e), "source": "Google Meet"}


async def meet_get_transcript(
    conference_id: str,
) -> dict[str, Any]:
    """Ottiene la trascrizione di una conferenza Google Meet.

    Args:
        conference_id: ID della conferenza

    Returns:
        dict con trascrizione (speaker + testo)
    """
    try:
        service = _get_workspace_service()
        if not service:
            return {"error": "Google Workspace not configured", "source": "Google Meet"}

        transcript = await service.meet_get_transcript(conference_id=conference_id)

        return {
            **transcript,
            "source": "Google Meet",
        }

    except Exception as e:
        logger.error("meet_get_transcript_error", error=str(e))
        return {"error": str(e), "source": "Google Meet"}


# =============================================================================
# Forms Tools
# =============================================================================


async def forms_get(form_id: str) -> dict[str, Any]:
    """Ottiene info su un Google Form.

    Args:
        form_id: ID del form

    Returns:
        dict con info form
    """
    try:
        service = _get_workspace_service()
        forms = await _run_sync(partial(service._get_google_service, "forms", "v1"))

        result = await _run_sync(lambda: forms.forms().get(formId=form_id).execute())

        items = result.get("items", [])

        return {
            "form_id": form_id,
            "title": result.get("info", {}).get("title", ""),
            "description": result.get("info", {}).get("description", ""),
            "item_count": len(items),
            "source": "Google Forms",
        }

    except Exception as e:
        logger.error("forms_get_error", error=str(e))
        return {"error": str(e), "source": "Google Forms"}


async def forms_get_responses(form_id: str) -> dict[str, Any]:
    """Ottiene risposte di un Google Form.

    Args:
        form_id: ID del form

    Returns:
        dict con risposte
    """
    try:
        service = _get_workspace_service()
        forms = await _run_sync(partial(service._get_google_service, "forms", "v1"))

        result = await _run_sync(lambda: forms.forms().responses().list(formId=form_id).execute())

        responses = result.get("responses", [])

        return {
            "form_id": form_id,
            "response_count": len(responses),
            "responses": responses[:20],  # Limita a 20
            "source": "Google Forms",
        }

    except Exception as e:
        logger.error("forms_get_responses_error", error=str(e))
        return {"error": str(e), "source": "Google Forms"}


# =============================================================================
# Classroom Tools
# =============================================================================


async def classroom_list_courses() -> dict[str, Any]:
    """Lista corsi Google Classroom.

    Returns:
        dict con lista corsi
    """
    try:
        service = _get_workspace_service()
        classroom = await _run_sync(partial(service._get_google_service, "classroom", "v1"))

        result = await _run_sync(lambda: classroom.courses().list(pageSize=50).execute())

        courses = result.get("courses", [])

        return {
            "courses": [
                {
                    "id": c.get("id"),
                    "name": c.get("name"),
                    "section": c.get("section", ""),
                    "state": c.get("courseState", ""),
                }
                for c in courses
            ],
            "count": len(courses),
            "source": "Google Classroom",
        }

    except Exception as e:
        logger.error("classroom_list_courses_error", error=str(e))
        return {"error": str(e), "source": "Google Classroom"}


async def classroom_get_coursework(course_id: str) -> dict[str, Any]:
    """Ottiene compiti di un corso Classroom.

    Args:
        course_id: ID del corso

    Returns:
        dict con lista compiti
    """
    try:
        service = _get_workspace_service()
        classroom = await _run_sync(partial(service._get_google_service, "classroom", "v1"))

        result = await _run_sync(
            lambda: (
                classroom.courses()
                .courseWork()
                .list(
                    courseId=course_id,
                )
                .execute()
            )
        )

        coursework = result.get("courseWork", [])

        return {
            "course_id": course_id,
            "coursework": [
                {
                    "id": cw.get("id"),
                    "title": cw.get("title"),
                    "state": cw.get("state", ""),
                    "dueDate": cw.get("dueDate"),
                }
                for cw in coursework
            ],
            "count": len(coursework),
            "source": "Google Classroom",
        }

    except Exception as e:
        logger.error("classroom_get_coursework_error", error=str(e))
        return {"error": str(e), "source": "Google Classroom"}


# =============================================================================
# Tool Registry
# =============================================================================

AVAILABLE_TOOLS = {
    # Drive (7 tools)
    "google_drive_search": drive_search,
    "google_drive_list_files": drive_list_files,
    "google_drive_get_file": drive_get_file,
    "google_drive_get_content": drive_get_content,  # NEW
    "google_drive_export": drive_export,  # NEW
    "google_drive_create_folder": drive_create_folder,  # NEW
    "google_drive_copy": drive_copy,  # NEW
    # Gmail (7 tools)
    "google_gmail_search": gmail_search,
    "google_gmail_get_message": gmail_get_message,
    "google_gmail_send": gmail_send,  # NEW
    "google_gmail_reply": gmail_reply,  # NEW
    "google_gmail_forward": gmail_forward,  # NEW
    "google_gmail_get_attachments": gmail_get_attachments,  # NEW
    "google_gmail_get_attachment_content": gmail_get_attachment_content,  # NEW
    # Calendar (7 tools)
    "google_calendar_upcoming": calendar_upcoming,
    "google_calendar_list_events": calendar_list_events,
    "google_calendar_create_event": calendar_create_event,  # NEW
    "google_calendar_get_event": calendar_get_event,  # NEW
    "google_calendar_update_event": calendar_update_event,  # NEW
    "google_calendar_delete_event": calendar_delete_event,  # NEW
    "google_calendar_analyze_meetings": calendar_analyze_meetings,  # NEW - Multi-platform meeting analyzer
    # Docs (5 tools)
    "google_docs_get": docs_get,
    "google_docs_create": docs_create,  # NEW
    "google_docs_insert_text": docs_insert_text,  # NEW
    "google_docs_append_text": docs_append_text,  # NEW
    "google_docs_replace_text": docs_replace_text,  # NEW
    # Sheets (6 tools)
    "google_sheets_get_values": sheets_get_values,
    "google_sheets_get_metadata": sheets_get_metadata,
    "google_sheets_create": sheets_create,  # NEW
    "google_sheets_update_values": sheets_update_values,  # NEW
    "google_sheets_append_row": sheets_append_row,  # NEW
    "google_sheets_add_sheet": sheets_add_sheet,  # NEW
    # Slides (4 tools)
    "google_slides_get": slides_get,
    "google_slides_create": slides_create,  # NEW
    "google_slides_get_text": slides_get_text,  # NEW
    "google_slides_add_slide": slides_add_slide,  # NEW
    # Meet (3 tools)
    "google_meet_create": meet_create,
    "google_meet_list_conferences": meet_list_conferences,  # NEW
    "google_meet_get_transcript": meet_get_transcript,  # NEW
    # Forms (2 tools)
    "google_forms_get": forms_get,
    "google_forms_get_responses": forms_get_responses,
    # Classroom (2 tools)
    "google_classroom_list_courses": classroom_list_courses,
    "google_classroom_get_coursework": classroom_get_coursework,
}


async def execute_tool(tool_name: str, arguments: dict[str, Any]) -> dict[str, Any]:
    """Esegue tool Google per nome.

    Filtra automaticamente parametri non accettati dalla funzione.

    Args:
        tool_name: Nome del tool
        arguments: Argomenti per il tool

    Returns:
        Risultato del tool
    """
    import inspect

    if tool_name not in AVAILABLE_TOOLS:
        return {
            "error": f"Unknown Google Workspace tool: {tool_name}",
            "available": list(AVAILABLE_TOOLS.keys()),
        }

    tool_func = AVAILABLE_TOOLS[tool_name]

    # Filter arguments to only those the function accepts
    sig = inspect.signature(tool_func)
    valid_params = set(sig.parameters.keys())
    filtered_args = {k: v for k, v in arguments.items() if k in valid_params}

    if len(filtered_args) < len(arguments):
        ignored = set(arguments.keys()) - valid_params
        logger.warning(
            "execute_tool_ignored_params",
            tool=tool_name,
            ignored=list(ignored),
            hint="LLM hallucinated parameters not in function signature",
        )

    return await tool_func(**filtered_args)


# =============================================================================
# Tool Engine Integration
# =============================================================================


def get_tool_definitions() -> list:
    """Generate ToolDefinition objects for all Google Workspace tools."""
    from me4brain.engine.types import ToolDefinition, ToolParameter

    definitions = [
        # Drive Tools
        ToolDefinition(
            name="google_drive_search",
            description=(
                "Search files on Google Drive. "
                "QUERY SYNTAX: 'fullText contains {keyword}' OR 'name contains {name}'. "
                "Combine with AND/OR. Use max_results=50-100 for complete search."
            ),
            parameters={
                "query": ToolParameter(
                    type="string",
                    description=(
                        "Search term. Examples: 'project report', 'ANCI', 'budget 2025'. "
                        "Full-text search in content and filename."
                    ),
                    required=True,
                ),
                "max_results": ToolParameter(
                    type="integer",
                    description="Results to return (1-100). Use 50+ for complete search.",
                    required=False,
                ),
                "mime_type": ToolParameter(
                    type="string",
                    description="Filter: 'application/pdf', 'application/vnd.google-apps.document'",
                    required=False,
                ),
                "folder_id": ToolParameter(
                    type="string",
                    description="Folder ID to restrict scope (from previous search)",
                    required=False,
                ),
            },
            domain="search",
            category="drive",
        ),
        ToolDefinition(
            name="google_drive_list_files",
            description="List files and folders in a Google Drive directory. Returns file names, types, sizes, and modification dates. Use when user asks to see contents of a folder, list their Drive files, or browse their cloud storage.",
            parameters={
                "folder_id": ToolParameter(
                    type="string",
                    description="Folder ID to list contents of (None = root/My Drive)",
                    required=False,
                ),
                "max_results": ToolParameter(
                    type="integer",
                    description="Maximum number of files to return (default 50)",
                    required=False,
                ),
            },
            domain="file_management",
            category="drive",
        ),
        ToolDefinition(
            name="google_drive_get_file",
            description="Get metadata and details of a specific file on Google Drive. Returns file name, size, owner, sharing permissions, and last modified date. Use when user asks for info about a specific file or needs file details before downloading/editing.",
            parameters={
                "file_id": ToolParameter(
                    type="string",
                    description="Google Drive file ID (from URL or previous search)",
                    required=True,
                ),
            },
            domain="file_management",
            category="drive",
        ),
        ToolDefinition(
            name="google_drive_get_content",
            description="Extract and read the full text content of a Google file (Docs, Sheets, Slides). Use when user asks to read, show, or get the content of a document, spreadsheet data, or presentation text.",
            parameters={
                "file_id": ToolParameter(
                    type="string",
                    description="Google Drive file ID of the document/sheet/slide",
                    required=True,
                ),
            },
            domain="file_management",
            category="drive",
        ),
        ToolDefinition(
            name="google_drive_export",
            description="Export a Google file to a standard format (PDF, DOCX, XLSX, PPTX, TXT, CSV). Converts Google Docs/Sheets/Slides to downloadable formats. Use when user asks to export, download, or convert a Google document to another format.",
            parameters={
                "file_id": ToolParameter(
                    type="string",
                    description="Google Drive file ID to export",
                    required=True,
                ),
                "export_format": ToolParameter(
                    type="string",
                    description="Output format: 'pdf', 'docx', 'xlsx', 'pptx', 'txt', 'csv'",
                    required=False,
                ),
            },
            domain="file_management",
            category="drive",
        ),
        ToolDefinition(
            name="google_drive_create_folder",
            description="Create a new folder on Google Drive. Use when user asks to create, make, or add a new folder for organizing files in their Drive.",
            parameters={
                "name": ToolParameter(
                    type="string", description="Name of the new folder", required=True
                ),
                "parent_id": ToolParameter(
                    type="string",
                    description="Parent folder ID (None = create in root/My Drive)",
                    required=False,
                ),
            },
            domain="file_management",
            category="drive",
        ),
        ToolDefinition(
            name="google_drive_copy",
            description="Copy an existing file on Google Drive to create a duplicate. Use when user asks to copy, duplicate, or clone a file, or create a backup of a document.",
            parameters={
                "file_id": ToolParameter(
                    type="string",
                    description="Google Drive file ID to copy",
                    required=True,
                ),
                "new_name": ToolParameter(
                    type="string",
                    description="Name for the new copy (default: 'Copy of original')",
                    required=False,
                ),
                "parent_id": ToolParameter(
                    type="string",
                    description="Destination folder ID (None = same folder as original)",
                    required=False,
                ),
            },
            domain="file_management",
            category="drive",
        ),
        # Gmail Tools
        ToolDefinition(
            name="google_gmail_search",
            description="Cerca nelle email, cerca nell'email. Trova nella posta Gmail informazioni, indicazioni, dettagli su eventi, appuntamenti, riunioni, progetti, viaggi, conferenze, meeting. Cerca nell'email informazioni relative a un evento. Trova nelle email comunicazioni su progetti, appuntamenti, trasferte. Usa questo strumento quando l'utente chiede: cerca nell'email, trova nella posta, cerca informazioni su un evento nell'email, trova dettagli di un progetto nelle email, cerca comunicazioni su un viaggio. Supports Gmail query: from:, to:, subject:, after:, before:, has:attachment.",
            parameters={
                "query": ToolParameter(
                    type="string",
                    description="Gmail search query (e.g., 'from:boss@company.com after:2024/01/01', 'subject:invoice has:attachment')",
                    required=True,
                ),
                "max_results": ToolParameter(
                    type="integer",
                    description="Maximum number of emails to return (default 10)",
                    required=False,
                ),
            },
            domain="search",
            category="gmail",
        ),
        ToolDefinition(
            name="google_gmail_get_message",
            description="Read the full content of a specific email including body, attachments info, and headers. Leggi email, apri messaggio, mostra posta. Use when user asks 'apri email', 'leggi messaggio', 'mostra email', 'open email', 'read message'.",
            parameters={
                "message_id": ToolParameter(
                    type="string",
                    description="Gmail message ID (from previous search results)",
                    required=True,
                ),
            },
            domain="communication",
            category="gmail",
        ),
        ToolDefinition(
            name="google_gmail_send",
            description="Compose and send a new email via Gmail. Invia email, manda messaggio, scrivi email. Supports CC, BCC, and plain text body. Use when user asks 'invia email', 'manda mail', 'scrivi messaggio', 'send email', 'compose mail'.",
            parameters={
                "to": ToolParameter(
                    type="string",
                    description="Recipient email address (e.g., 'recipient@example.com')",
                    required=True,
                ),
                "subject": ToolParameter(
                    type="string", description="Email subject line", required=True
                ),
                "body": ToolParameter(
                    type="string",
                    description="Email body content (plain text)",
                    required=True,
                ),
                "cc": ToolParameter(
                    type="string",
                    description="CC recipient(s), comma-separated",
                    required=False,
                ),
                "bcc": ToolParameter(
                    type="string",
                    description="BCC recipient(s), comma-separated",
                    required=False,
                ),
            },
            domain="communication",
            category="gmail",
        ),
        ToolDefinition(
            name="google_gmail_reply",
            description="Reply to an existing email thread, keeping the conversation context. Rispondi a email, rispondi al messaggio. Use when user asks 'rispondi all'email', 'rispondi al messaggio', 'reply to email', 'respond to message'.",
            parameters={
                "message_id": ToolParameter(
                    type="string",
                    description="Gmail message ID to reply to (from previous search)",
                    required=True,
                ),
                "body": ToolParameter(
                    type="string", description="Reply message content", required=True
                ),
            },
            domain="communication",
            category="gmail",
        ),
        ToolDefinition(
            name="google_gmail_forward",
            description="Forward an email to another recipient with optional additional message. Use when user asks to forward, share, or send an email to someone else.",
            parameters={
                "message_id": ToolParameter(
                    type="string",
                    description="Gmail message ID to forward",
                    required=True,
                ),
                "to": ToolParameter(
                    type="string",
                    description="Recipient email address to forward to",
                    required=True,
                ),
                "additional_text": ToolParameter(
                    type="string",
                    description="Optional message to include with forwarded email",
                    required=False,
                ),
            },
            domain="communication",
            category="gmail",
        ),
        ToolDefinition(
            name="google_gmail_get_attachments",
            description="Get attachments from a specific email. Returns list of attachment filenames, types and sizes. Use when user asks about email attachments, what files were attached, or needs to analyze attachment content.",
            parameters={
                "message_id": ToolParameter(
                    type="string",
                    description="Gmail message ID to get attachments from",
                    required=True,
                ),
            },
            domain="communication",
            category="gmail",
        ),
        ToolDefinition(
            name="google_gmail_get_attachment_content",
            description="Download and extract text content from an email attachment. Supports PDF, Word, PowerPoint, Excel, images (OCR) and text files. Use when user asks to read, analyze, or extract content from email attachments.",
            parameters={
                "message_id": ToolParameter(
                    type="string",
                    description="Gmail message ID containing the attachment",
                    required=True,
                ),
                "attachment_id": ToolParameter(
                    type="string",
                    description="Attachment ID from gmail_get_attachments response",
                    required=True,
                ),
            },
            domain="communication",
            category="gmail",
        ),
        # Calendar Tools
        ToolDefinition(
            name="google_calendar_upcoming",
            description="Get upcoming calendar events for the next N days. Returns meeting titles, times, locations, and attendees. Use when user asks 'what's on my calendar', 'my schedule for this week', 'upcoming meetings', or 'what do I have tomorrow'.",
            parameters={
                "days": ToolParameter(
                    type="integer",
                    description="Number of days to look ahead (default 7)",
                    required=False,
                ),
                "calendar_id": ToolParameter(
                    type="string",
                    description="Calendar ID (default: primary/main calendar)",
                    required=False,
                ),
            },
            domain="scheduling",
            category="calendar",
        ),
        ToolDefinition(
            name="google_calendar_list_events",
            description="Search and list calendar events with date range filters. Find events by title or within a specific time period. Use when user asks to find a specific meeting, list events between dates, or search calendar history.",
            parameters={
                "query": ToolParameter(
                    type="string",
                    description="Search term to find in event titles",
                    required=False,
                ),
                "time_min": ToolParameter(
                    type="string",
                    description="Start date/time filter (ISO format, e.g., '2024-01-01T00:00:00Z')",
                    required=False,
                ),
                "time_max": ToolParameter(
                    type="string",
                    description="End date/time filter (ISO format)",
                    required=False,
                ),
                "max_results": ToolParameter(
                    type="integer",
                    description="Maximum events to return (default 50)",
                    required=False,
                ),
            },
            domain="scheduling",
            category="calendar",
        ),
        ToolDefinition(
            name="google_calendar_create_event",
            description="Create a new event on Google Calendar. Schedule meetings, appointments, or reminders with title, time, description, and location. Use when user asks to schedule, create, add, or book a meeting or event.",
            parameters={
                "summary": ToolParameter(
                    type="string", description="Event title/name", required=True
                ),
                "start_time": ToolParameter(
                    type="string",
                    description="Start date/time (ISO format, e.g., '2024-01-15T14:00:00')",
                    required=True,
                ),
                "end_time": ToolParameter(
                    type="string",
                    description="End date/time (ISO format)",
                    required=True,
                ),
                "description": ToolParameter(
                    type="string",
                    description="Event description or notes",
                    required=False,
                ),
                "location": ToolParameter(
                    type="string",
                    description="Event location or meeting room",
                    required=False,
                ),
            },
            domain="scheduling",
            category="calendar",
        ),
        ToolDefinition(
            name="google_calendar_get_event",
            description="Get full details of a specific calendar event including attendees, attachments, and conferencing info. Use when user asks for details about a particular meeting or event.",
            parameters={
                "event_id": ToolParameter(
                    type="string",
                    description="Calendar event ID (from list/search results)",
                    required=True,
                ),
                "calendar_id": ToolParameter(
                    type="string",
                    description="Calendar ID (default: primary)",
                    required=False,
                ),
            },
            domain="scheduling",
            category="calendar",
        ),
        ToolDefinition(
            name="google_calendar_update_event",
            description="Update an existing calendar event. Modify event title, time, description, or location. Use when user asks to reschedule, change, update, or move a meeting.",
            parameters={
                "event_id": ToolParameter(
                    type="string",
                    description="Calendar event ID to update",
                    required=True,
                ),
                "summary": ToolParameter(
                    type="string",
                    description="New event title (optional)",
                    required=False,
                ),
                "start_time": ToolParameter(
                    type="string",
                    description="New start time (ISO format, optional)",
                    required=False,
                ),
                "end_time": ToolParameter(
                    type="string",
                    description="New end time (ISO format, optional)",
                    required=False,
                ),
            },
            domain="scheduling",
            category="calendar",
        ),
        ToolDefinition(
            name="google_calendar_delete_event",
            description="Delete an event from Google Calendar. Permanently removes the meeting or appointment. Use when user asks to cancel, delete, or remove a meeting.",
            parameters={
                "event_id": ToolParameter(
                    type="string",
                    description="Calendar event ID to delete",
                    required=True,
                ),
                "calendar_id": ToolParameter(
                    type="string",
                    description="Calendar ID (default: primary)",
                    required=False,
                ),
            },
            domain="scheduling",
            category="calendar",
        ),
        # Docs Tools
        ToolDefinition(
            name="google_docs_get",
            description="Read the full text content of a Google Docs document. Returns document body, headers, and structure. Use when user asks to read, show, open, or get content of a Google Doc.",
            parameters={
                "document_id": ToolParameter(
                    type="string",
                    description="Google Docs document ID (from Drive or URL)",
                    required=True,
                ),
            },
            domain="content_creation",
            category="docs",
        ),
        ToolDefinition(
            name="google_docs_create",
            description="Create a new Google Docs document with optional initial content and save it to a specific folder. Use when user asks to create, make, or start a new document, report, or text file. Also use when user wants to TRANSFORM, CONVERT, EXPORT, or SAVE a response, text, or content AS a Google Doc document. Can optionally save directly into a specified Google Drive folder.",
            parameters={
                "title": ToolParameter(
                    type="string", description="Document title/name", required=True
                ),
                "content": ToolParameter(
                    type="string",
                    description="Initial document content (optional)",
                    required=False,
                ),
                "folder_id": ToolParameter(
                    type="string",
                    description="Google Drive folder ID where to save the document (optional, default: root). Extract from folder URL like https://drive.google.com/drive/folders/{folder_id}",
                    required=False,
                ),
            },
            domain="content_creation",
            category="docs",
        ),
        ToolDefinition(
            name="google_docs_insert_text",
            description="Insert text at a specific position in a Google Docs document. Use when user asks to add text at a particular location in a document.",
            parameters={
                "document_id": ToolParameter(
                    type="string", description="Google Docs document ID", required=True
                ),
                "text": ToolParameter(
                    type="string", description="Text content to insert", required=True
                ),
                "index": ToolParameter(
                    type="integer",
                    description="Character position to insert at (default: beginning)",
                    required=False,
                ),
            },
            domain="content_creation",
            category="docs",
        ),
        ToolDefinition(
            name="google_docs_append_text",
            description="Append text to the end of a Google Docs document. Use when user asks to add content to the end of a document without overwriting.",
            parameters={
                "document_id": ToolParameter(
                    type="string", description="Google Docs document ID", required=True
                ),
                "text": ToolParameter(
                    type="string",
                    description="Text to append at the end",
                    required=True,
                ),
            },
            domain="content_creation",
            category="docs",
        ),
        ToolDefinition(
            name="google_docs_replace_text",
            description="Find and replace text in a Google Docs document. Performs global search-and-replace. Use when user asks to replace, substitute, or change specific words or phrases throughout a document.",
            parameters={
                "document_id": ToolParameter(
                    type="string", description="Google Docs document ID", required=True
                ),
                "search_text": ToolParameter(
                    type="string",
                    description="Text to find (exact match)",
                    required=True,
                ),
                "replace_text": ToolParameter(
                    type="string", description="Replacement text", required=True
                ),
            },
            domain="content_creation",
            category="docs",
        ),
        # Sheets Tools
        ToolDefinition(
            name="google_sheets_get_values",
            description="Read cell values from a Google Sheets spreadsheet. Returns data from a specified range (e.g., 'A1:D10', 'Sheet1!A:A'). Use when user asks to read, get, or show spreadsheet data.",
            parameters={
                "spreadsheet_id": ToolParameter(
                    type="string",
                    description="Google Sheets spreadsheet ID (from URL)",
                    required=True,
                ),
                "range": ToolParameter(
                    type="string",
                    description="Cell range in A1 notation (e.g., 'A1:D10', 'Sheet1!A1:B50')",
                    required=True,
                ),
            },
            domain="data_analysis",
            category="sheets",
        ),
        ToolDefinition(
            name="google_sheets_get_metadata",
            description="Get metadata and properties of a Google Sheets spreadsheet. Returns title, sheet names, column/row counts. Use when user asks about spreadsheet structure or details.",
            parameters={
                "spreadsheet_id": ToolParameter(
                    type="string",
                    description="Google Sheets spreadsheet ID",
                    required=True,
                ),
            },
            domain="data_analysis",
            category="sheets",
        ),
        ToolDefinition(
            name="google_sheets_create",
            description="Create a new Google Sheets spreadsheet. Use when user asks to create, make, or start a new spreadsheet, workbook, or Excel-like file.",
            parameters={
                "title": ToolParameter(
                    type="string", description="Spreadsheet title/name", required=True
                ),
            },
            domain="data_analysis",
            category="sheets",
        ),
        ToolDefinition(
            name="google_sheets_update_values",
            description="Update cell values in a Google Sheets spreadsheet. Writes data to a specified range. Use when user asks to update, change, modify, or write data to a spreadsheet.",
            parameters={
                "spreadsheet_id": ToolParameter(
                    type="string",
                    description="Google Sheets spreadsheet ID",
                    required=True,
                ),
                "range": ToolParameter(
                    type="string",
                    description="Target cell range (e.g., 'A1:C10')",
                    required=True,
                ),
                "values": ToolParameter(
                    type="array",
                    description="2D array of values to write (e.g., [['A1', 'B1'], ['A2', 'B2']])",
                    required=True,
                ),
            },
            domain="data_analysis",
            category="sheets",
        ),
        ToolDefinition(
            name="google_sheets_append_row",
            description="Append a new row to the end of a Google Sheets spreadsheet. Adds data after the last row with content. Use when user asks to add a row, insert data at the end, or log new entries.",
            parameters={
                "spreadsheet_id": ToolParameter(
                    type="string",
                    description="Google Sheets spreadsheet ID",
                    required=True,
                ),
                "range": ToolParameter(
                    type="string",
                    description="Sheet name or range (e.g., 'Sheet1' or 'Sheet1!A:D')",
                    required=True,
                ),
                "values": ToolParameter(
                    type="array",
                    description="Row values as array (e.g., ['value1', 'value2', 'value3'])",
                    required=True,
                ),
            },
            domain="data_analysis",
            category="sheets",
        ),
        ToolDefinition(
            name="google_sheets_add_sheet",
            description="Add a new sheet (tab) to an existing Google Sheets spreadsheet. Use when user asks to create a new tab, add a worksheet, or create another sheet within a spreadsheet.",
            parameters={
                "spreadsheet_id": ToolParameter(
                    type="string",
                    description="Google Sheets spreadsheet ID",
                    required=True,
                ),
                "title": ToolParameter(
                    type="string",
                    description="Name for the new sheet/tab",
                    required=True,
                ),
            },
            domain="data_analysis",
            category="sheets",
        ),
        # Slides Tools
        ToolDefinition(
            name="google_slides_get",
            description="Get information and metadata about a Google Slides presentation. Returns title, slide count, and structure. Use when user asks for presentation details or properties.",
            parameters={
                "presentation_id": ToolParameter(
                    type="string",
                    description="Google Slides presentation ID (from URL)",
                    required=True,
                ),
            },
            domain="content_creation",
            category="slides",
        ),
        ToolDefinition(
            name="google_slides_create",
            description="Create a new Google Slides presentation. Use when user asks to create, make, or start a new presentation, slideshow, or pitch deck.",
            parameters={
                "title": ToolParameter(
                    type="string", description="Presentation title", required=True
                ),
            },
            domain="content_creation",
            category="slides",
        ),
        ToolDefinition(
            name="google_slides_get_text",
            description="Extract all text content from a Google Slides presentation. Returns text from all slides. Use when user asks to read, get content, or extract text from a presentation.",
            parameters={
                "presentation_id": ToolParameter(
                    type="string",
                    description="Google Slides presentation ID",
                    required=True,
                ),
            },
            domain="content_creation",
            category="slides",
        ),
        ToolDefinition(
            name="google_slides_add_slide",
            description="Add a new slide to a Google Slides presentation. Supports different layouts (blank, title, title+body). Use when user asks to add a slide, insert a new page, or expand a presentation.",
            parameters={
                "presentation_id": ToolParameter(
                    type="string",
                    description="Google Slides presentation ID",
                    required=True,
                ),
                "layout": ToolParameter(
                    type="string",
                    description="Slide layout: 'BLANK', 'TITLE', 'TITLE_AND_BODY', 'TITLE_ONLY'",
                    required=False,
                ),
            },
            domain="content_creation",
            category="slides",
        ),
        # Meet Tool
        ToolDefinition(
            name="google_meet_create",
            description="Create a new Google Meet video conference with a shareable link. Schedules a meeting with video conferencing enabled. Use when user asks to create a meeting, start a video call, or schedule a Google Meet.",
            parameters={
                "summary": ToolParameter(
                    type="string", description="Meeting title/subject", required=True
                ),
                "start_time": ToolParameter(
                    type="string",
                    description="Meeting start time (ISO format, e.g., '2024-01-15T14:00:00')",
                    required=True,
                ),
                "duration_minutes": ToolParameter(
                    type="integer",
                    description="Meeting duration in minutes (default: 60)",
                    required=False,
                ),
            },
            domain="scheduling",
            category="meet",
        ),
        ToolDefinition(
            name="google_meet_list_conferences",
            description="List past Google Meet conferences with participants. Shows meeting history with start/end times and attendee list. Use when user asks 'show my past meetings', 'who attended the meeting', 'list meet history', or needs to find a specific past call.",
            parameters={
                "max_results": ToolParameter(
                    type="integer",
                    description="Maximum number of conferences to return (default: 50)",
                    required=False,
                ),
            },
            domain="scheduling",
            category="meet",
        ),
        ToolDefinition(
            name="google_meet_get_transcript",
            description="Get the transcript of a Google Meet conference. Returns speaker names and what they said during the meeting. Use when user asks for 'meeting transcript', 'what was said in the meeting', 'meeting notes', or 'conversation from the call'.",
            parameters={
                "conference_id": ToolParameter(
                    type="string",
                    description="Google Meet conference ID (from list_conferences)",
                    required=True,
                ),
            },
            domain="content_creation",
            category="meet",
        ),
        # Forms Tools
        ToolDefinition(
            name="google_forms_get",
            description="Get information about a Google Forms form including title, description, and question count. Use when user asks for details about a form or survey.",
            parameters={
                "form_id": ToolParameter(
                    type="string",
                    description="Google Forms form ID (from URL)",
                    required=True,
                ),
            },
            domain="content_creation",
            category="forms",
        ),
        ToolDefinition(
            name="google_forms_get_responses",
            description="Get all responses submitted to a Google Forms form. Returns individual responses with answers. Use when user asks for form results, survey responses, or submitted data.",
            parameters={
                "form_id": ToolParameter(
                    type="string", description="Google Forms form ID", required=True
                ),
            },
            domain="data_analysis",
            category="forms",
        ),
        # Classroom Tools
        ToolDefinition(
            name="google_classroom_list_courses",
            description="List all Google Classroom courses the user is enrolled in or teaches. Returns course names, IDs, and states. Use when user asks about their courses, classes, or Classroom.",
            parameters={},
            domain="communication",
            category="classroom",
        ),
        ToolDefinition(
            name="google_classroom_get_coursework",
            description="Get assignments and coursework from a Google Classroom course. Returns assignment titles, due dates, and states. Use when user asks about homework, assignments, or course tasks.",
            parameters={
                "course_id": ToolParameter(
                    type="string",
                    description="Google Classroom course ID (from list_courses)",
                    required=True,
                ),
            },
            domain="communication",
            category="classroom",
        ),
    ]

    return definitions


def get_executors() -> dict:
    """Return mapping of tool names to executor functions."""
    return AVAILABLE_TOOLS
